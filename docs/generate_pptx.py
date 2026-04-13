"""
Generate the OTA Firmware Update System presentation as .pptx
Run: python docs/generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ──────────────────────────────────────────
BG_DARK   = RGBColor(0x0E, 0x11, 0x16)
BG_CARD   = RGBColor(0x16, 0x1B, 0x22)
ACCENT    = RGBColor(0x58, 0xA6, 0xFF)
GREEN     = RGBColor(0x7A, 0xDF, 0x7A)
RED       = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x8B, 0x94, 0x9E)
LIGHT     = RGBColor(0xD6, 0xDD, 0xE6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def _set_bg(slide, colour=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def _add_shape(slide, left, top, width, height, fill_colour=BG_CARD, corner=Inches(0.15)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_colour
    shp.line.fill.background()
    shp.rotation = 0.0
    return shp

def _text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def _set_text(tf, text, size=18, colour=LIGHT, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = colour
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align

def _mono(tf, text, size=14, colour=LIGHT):
    _set_text(tf, text, size=size, colour=colour, font_name="Consolas")

def _title_subtitle(slide, title, subtitle=""):
    _set_bg(slide)
    tb = _text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    _set_text(tb.text_frame, title, size=32, colour=WHITE, bold=True)
    if subtitle:
        tb2 = _text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
        _set_text(tb2.text_frame, subtitle, size=16, colour=GREY)

# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
# Accent bar
_add_shape(sl, 0, 0, W, Inches(0.08), ACCENT)
tb = _text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(1.5))
_set_text(tb.text_frame, "S-Board OTA Firmware Update System", size=44, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb2 = _text_box(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6))
_set_text(tb2.text_frame, "CAN-FD Over-The-Air Updates for TI C2000 F28P55x", size=22, colour=ACCENT, align=PP_ALIGN.CENTER)
tb3 = _text_box(sl, Inches(1), Inches(4.8), Inches(11), Inches(0.5))
_set_text(tb3.text_frame, "GEN3 Power Distribution Unit  |  Delta Electronics", size=16, colour=GREY, align=PP_ALIGN.CENTER)
tb4 = _text_box(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.5))
_set_text(tb4.text_frame, "Author: Parthasarathy M", size=14, colour=GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "The Problem")
card = _add_shape(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
tb = _text_box(sl, Inches(1.2), Inches(2.0), Inches(4.8), Inches(4.6))
lines = [
    ("Updating 100+ deployed S-Boards:", LIGHT, 20, True),
    ("", LIGHT, 10, False),
    ("   Physical access to each board", RED, 18, False),
    ("   JTAG debugger connection", RED, 18, False),
    ("   CCS installed on field laptop", RED, 18, False),
    ("   15+ minutes per board", RED, 18, False),
    ("   Risk of bricking on power loss", RED, 18, False),
]
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, col, sz, bld) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = col; p.font.bold = bld; p.font.name = "Segoe UI"

card2 = _add_shape(sl, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), RGBColor(0x0D, 0x1F, 0x0D))
tb2 = _text_box(sl, Inches(7.2), Inches(2.0), Inches(4.8), Inches(4.6))
lines2 = [
    ("What we need:", GREEN, 20, True),
    ("", LIGHT, 10, False),
    ("   Remote update over CAN bus", GREEN, 18, False),
    ("   No JTAG, no CCS needed", GREEN, 18, False),
    ("   ~30 seconds per board", GREEN, 18, False),
    ("   Brick-proof design", GREEN, 18, False),
    ("   Professional GUI tool", GREEN, 18, False),
]
tf2 = tb2.text_frame; tf2.word_wrap = True
for i, (txt, col, sz, bld) in enumerate(lines2):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = col; p.font.bold = bld; p.font.name = "Segoe UI"

# ════════════════════════════════════════════════════════════════
# SLIDE 3: Dual-Bank Architecture
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Dual-Bank OTA Architecture", "Flash memory layout on F28P55x")

# Bank 0 box
b0 = _add_shape(sl, Inches(1), Inches(2), Inches(4.5), Inches(4.5), RGBColor(0x1A, 0x2A, 0x1A))
tb = _text_box(sl, Inches(1.3), Inches(2.1), Inches(4), Inches(0.5))
_set_text(tb.text_frame, "BANK 0  (0x080000)", size=18, colour=GREEN, bold=True)
# Boot mgr section
bm = _add_shape(sl, Inches(1.3), Inches(2.7), Inches(3.8), Inches(1.2), RGBColor(0x2A, 0x1A, 0x1A))
tb = _text_box(sl, Inches(1.5), Inches(2.8), Inches(3.4), Inches(1))
_set_text(tb.text_frame, "Sectors 0-7: BOOT MANAGER\nFlashed via JTAG once\nNEVER erased by OTA", size=14, colour=YELLOW)
# App section
ap = _add_shape(sl, Inches(1.3), Inches(4.1), Inches(3.8), Inches(2.2), RGBColor(0x1A, 0x1A, 0x2A))
tb = _text_box(sl, Inches(1.5), Inches(4.2), Inches(3.4), Inches(2))
_set_text(tb.text_frame, "Sectors 8-127: APPLICATION\nRunning firmware\nUpdated by boot manager\nvia Bank 2 copy", size=14, colour=ACCENT)

# Bank 2 box
b2 = _add_shape(sl, Inches(7), Inches(2), Inches(4.5), Inches(4.5), RGBColor(0x1A, 0x1A, 0x2A))
tb = _text_box(sl, Inches(7.3), Inches(2.1), Inches(4), Inches(0.5))
_set_text(tb.text_frame, "BANK 2  (0x0C0000)", size=18, colour=ACCENT, bold=True)
stg = _add_shape(sl, Inches(7.3), Inches(2.7), Inches(3.8), Inches(3.6), RGBColor(0x0E, 0x1E, 0x3E))
tb = _text_box(sl, Inches(7.5), Inches(3.5), Inches(3.4), Inches(2))
_set_text(tb.text_frame, "128 sectors\nOTA STAGING AREA\n\nNew firmware received\nover CAN-FD\nwritten here first", size=16, colour=LIGHT)

# Arrow: Bank 2 → Bank 0
arrow = slide = sl  # draw text arrow
tb = _text_box(sl, Inches(5.6), Inches(4), Inches(1.4), Inches(0.6))
_set_text(tb.text_frame, "COPY", size=20, colour=GREEN, bold=True, align=PP_ALIGN.CENTER)
tb = _text_box(sl, Inches(5.6), Inches(4.5), Inches(1.4), Inches(0.5))
_set_text(tb.text_frame, "<<<<<<", size=24, colour=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Bank 3 small box
b3 = _add_shape(sl, Inches(7), Inches(6.6), Inches(4.5), Inches(0.7), RGBColor(0x2A, 0x2A, 0x1A))
tb = _text_box(sl, Inches(7.2), Inches(6.65), Inches(4), Inches(0.6))
_set_text(tb.text_frame, "BANK 3 (0x0E0000) — Boot Flag: 0xA5A5 / 0x5A5A + size + CRC", size=12, colour=YELLOW)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: Protocol
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "OTA Protocol — 4 Steps")

steps = [
    ("1. ERASE",  "CMD_FW_START (0x30)\nErase Bank 2 (128 sectors)\nWait ~10s for ACK", ACCENT),
    ("2. HEADER", "CMD_FW_HEADER (0x31)\nSend size, CRC, version\nMCU parses and ACKs", ACCENT),
    ("3. DATA",   "Data frames (ID 6)\n64 bytes each, 16/burst\nACK after each burst", GREEN),
    ("4. VERIFY", "CMD_FW_COMPLETE (0x33)\nMCU computes CRC32\nCRC_PASS or CRC_FAIL", YELLOW),
]
for i, (title, desc, colour) in enumerate(steps):
    x = Inches(0.6 + i * 3.1)
    card = _add_shape(sl, x, Inches(2), Inches(2.9), Inches(3.5), BG_CARD)
    # Step number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), Inches(2.2), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = colour; circ.line.fill.background()
    tb = _text_box(sl, x + Inches(1.05), Inches(2.25), Inches(0.7), Inches(0.65))
    _set_text(tb.text_frame, str(i+1), size=28, colour=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    # Title
    tb = _text_box(sl, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.5))
    _set_text(tb.text_frame, title, size=20, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Description
    tb = _text_box(sl, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(1.5))
    _set_text(tb.text_frame, desc, size=14, colour=LIGHT, align=PP_ALIGN.CENTER)

# Bottom bar
tb = _text_box(sl, Inches(0.8), Inches(6), Inches(11), Inches(0.8))
_set_text(tb.text_frame,
    "On CRC_PASS: MCU writes boot flag to Bank 3, resets.  Boot manager copies Bank 2 → Bank 0 on next boot.",
    size=16, colour=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: Brick-Proof
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Brick-Proof Design", "Every failure mode has a safe recovery path")

failures = [
    ("Power lost during Bank 2 write",    "Boot manager untouched. Boots old FW. Retry OTA."),
    ("Power lost during Bank 0 copy",     "Boot manager retries copy on next boot (flag still set)."),
    ("CRC mismatch after transfer",       "Boot manager clears flag, boots old FW. No harm."),
    ("CAN bus disconnected mid-transfer", "Sender times out. MCU stays on old FW. Retry later."),
]
for i, (scenario, recovery) in enumerate(failures):
    y = Inches(2.0 + i * 1.25)
    # Red scenario card
    _add_shape(sl, Inches(0.8), y, Inches(5.2), Inches(1), RGBColor(0x2A, 0x12, 0x12))
    tb = _text_box(sl, Inches(1.0), y + Inches(0.15), Inches(4.8), Inches(0.7))
    _set_text(tb.text_frame, scenario, size=16, colour=RED, bold=True)
    # Green recovery card
    _add_shape(sl, Inches(6.5), y, Inches(5.8), Inches(1), RGBColor(0x12, 0x2A, 0x12))
    tb = _text_box(sl, Inches(6.7), y + Inches(0.15), Inches(5.4), Inches(0.7))
    _set_text(tb.text_frame, recovery, size=15, colour=GREEN)

tb = _text_box(sl, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "Boot manager NEVER erases itself. Only a JTAG write to sectors 0-7 can brick the device.",
          size=16, colour=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: GUI
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Professional GUI — 8 Tabs")

tabs = [
    ("Manual Control",  "Step-by-step OTA\nConnect, Erase, Header\nData, Verify, Single Frame"),
    ("MCU Operations",  "Read/Write boot flag\nRead flash, Erase sector\nCRC, Reset, Raw CAN"),
    ("CAN Monitor",     "Live bus sniffer\nAll IDs, timestamps\nStart/Stop/Clear"),
    ("CAN Bus",         "PCAN channel select\nBit timing config\nLive bitrate display"),
    ("Protocol",        "CAN IDs (6, 7, 8)\nCommand codes\nResponse codes"),
    ("Timing",          "Burst size, frame size\nInter-frame delay\nAll timeouts, retries"),
    ("Header",          "Magic, image type\nDest bank, entry point\nFor custom flash layouts"),
    ("Settings",        "Dark/Light/System theme\nSave/Load settings\nReset to defaults"),
]
for i, (name, desc) in enumerate(tabs):
    col = i % 4; row = i // 4
    x = Inches(0.6 + col * 3.1); y = Inches(1.8 + row * 2.7)
    card = _add_shape(sl, x, y, Inches(2.9), Inches(2.4), BG_CARD)
    tb = _text_box(sl, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.5))
    _set_text(tb.text_frame, name, size=17, colour=ACCENT, bold=True)
    tb = _text_box(sl, x + Inches(0.2), y + Inches(0.8), Inches(2.5), Inches(1.4))
    _set_text(tb.text_frame, desc, size=13, colour=LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: Binary Fix
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Binary Format Fix", "hex2000 --image with ROMS directive")

# BROKEN side
_add_shape(sl, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5), RGBColor(0x2A, 0x12, 0x12))
tb = _text_box(sl, Inches(0.8), Inches(1.9), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "BROKEN  (hex2000 --binary)", size=20, colour=RED, bold=True)
tb = _text_box(sl, Inches(0.8), Inches(2.6), Inches(5), Inches(3.8))
_mono(tb.text_frame,
    "Byte 0:  codestart[0]     OK\n"
    "Byte 2:  codestart[1]     OK\n"
    "Byte 4:  FPUmathTables    WRONG!\n"
    "         (should be at byte 16)\n"
    "         Sections concatenated\n"
    "         with NO gaps\n"
    "\n"
    "Every section after codestart\n"
    "lands at the WRONG address.\n"
    "Firmware is SCRAMBLED.\n"
    "\n"
    "Size: 40,442 bytes (34 short)",
    size=13, colour=LIGHT)

# FIXED side
_add_shape(sl, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), RGBColor(0x0D, 0x1F, 0x0D))
tb = _text_box(sl, Inches(7.0), Inches(1.9), Inches(5), Inches(0.5))
_set_text(tb.text_frame, "FIXED  (hex2000 --image + ROMS)", size=20, colour=GREEN, bold=True)
tb = _text_box(sl, Inches(7.0), Inches(2.6), Inches(5), Inches(3.8))
_mono(tb.text_frame,
    "Byte 0:  codestart[0]     OK\n"
    "Byte 2:  codestart[1]     OK\n"
    "Byte 4:  0xFFFF (fill)    GAP\n"
    "Byte 14: 0xFFFF (fill)    GAP\n"
    "Byte 16: FPUmathTables    CORRECT!\n"
    "         Address gaps filled\n"
    "         with 0xFFFF\n"
    "\n"
    "Every section at the correct\n"
    "flash address. Works perfectly.\n"
    "\n"
    "Size: 40,476 bytes (exact)",
    size=13, colour=LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: Bugs Fixed
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Three Critical Bugs Fixed")

bugs = [
    ("BUG 1: Flash ops ran from FLASH",
     "boot_manager.c: eraseSector(), programEightWords() in .text (flash).\n"
     "Flash FSM locks ALL banks during erase — CPU stalls.",
     "FIX: #pragma CODE_SECTION(.TI.ramfunc) on all 5 functions.",
     RED, GREEN),
    ("BUG 2: Boot flag write failed",
     "fw_image_rx.c: FW_triggerUpdate() ran from flash, called\n"
     "FW_writeBootFlag() (RAM). Return address in locked flash.",
     "FIX: Both functions .TI.ramfunc. Self-contained Flash API init.",
     RED, GREEN),
    ("BUG 3: Binary format wrong",
     "hex2000 --binary concatenates sections, destroying address gaps.\n"
     "Firmware scrambled when copied to Bank 0.",
     "FIX: hex_image.hexcmd with --image + ROMS directive.",
     RED, GREEN),
]
for i, (title, desc, fix, tcol, fcol) in enumerate(bugs):
    y = Inches(1.8 + i * 1.8)
    _add_shape(sl, Inches(0.6), y, Inches(11.5), Inches(1.6), BG_CARD)
    tb = _text_box(sl, Inches(0.9), y + Inches(0.1), Inches(11), Inches(0.4))
    _set_text(tb.text_frame, title, size=18, colour=tcol, bold=True)
    tb = _text_box(sl, Inches(0.9), y + Inches(0.55), Inches(6.5), Inches(0.8))
    _set_text(tb.text_frame, desc, size=12, colour=LIGHT)
    tb = _text_box(sl, Inches(7.5), y + Inches(0.55), Inches(4.5), Inches(0.8))
    _set_text(tb.text_frame, fix, size=13, colour=fcol, bold=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Performance
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Performance")

_add_shape(sl, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5), BG_CARD)
tb = _text_box(sl, Inches(0.9), Inches(2.0), Inches(5), Inches(4.5))
specs = (
    "Specification\n"
    "─────────────────────────\n"
    "Firmware size:      ~40 KB\n"
    "CAN-FD data rate:   2 Mbps\n"
    "Frame size:         64 bytes\n"
    "Total frames:       ~633\n"
    "Burst size:         16 frames\n"
    "Inter-frame delay:  1 ms\n"
    "CRC algorithm:      CRC32 (IEEE 802.3)"
)
_mono(tb.text_frame, specs, size=15, colour=LIGHT)

_add_shape(sl, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), RGBColor(0x0D, 0x1F, 0x0D))
tb = _text_box(sl, Inches(7.1), Inches(2.0), Inches(5), Inches(4.5))
timing = (
    "Timing Breakdown\n"
    "─────────────────────────\n"
    "Bank 2 erase:       ~10 sec\n"
    "Data transfer:      ~15 sec\n"
    "CRC verification:    ~2 sec\n"
    "Boot manager copy:   ~5 sec\n"
    "─────────────────────────\n"
    "TOTAL:              ~32 sec\n"
    "\n"
    "vs JTAG: 15+ minutes"
)
_mono(tb.text_frame, timing, size=15, colour=GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: Repos
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_title_subtitle(sl, "Repository Structure")

repos = [
    ("S-Board-Firmware",  "github.com/partham20", "Application firmware\nOTA receiver (fw_image_rx.c)\nhex_image.hexcmd", ACCENT),
    ("boot_manager",      "github.com/partham20", "Boot manager\nBank 2 → Bank 0 copy\nCAN debug telemetry", YELLOW),
    ("fw-updater-gui",    "github.com/partham20", "Python GUI + CLI sender\nDocumentation\nPresentation", GREEN),
    ("powercalculation",  "gitlab.deltaww.com",   "Primary S-Board repo\nBranch: FW_OTA\nInternal mirror", GREY),
]
for i, (name, host, desc, colour) in enumerate(repos):
    y = Inches(1.8 + i * 1.3)
    _add_shape(sl, Inches(0.6), y, Inches(11.5), Inches(1.1), BG_CARD)
    tb = _text_box(sl, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.4))
    _set_text(tb.text_frame, name, size=20, colour=colour, bold=True)
    tb = _text_box(sl, Inches(0.9), y + Inches(0.55), Inches(3.5), Inches(0.4))
    _set_text(tb.text_frame, host, size=12, colour=GREY)
    tb = _text_box(sl, Inches(5), y + Inches(0.15), Inches(7), Inches(0.8))
    _set_text(tb.text_frame, desc, size=13, colour=LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: Thank You
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_shape(sl, 0, 0, W, Inches(0.08), ACCENT)
tb = _text_box(sl, Inches(1), Inches(2.5), Inches(11), Inches(1))
_set_text(tb.text_frame, "Thank You", size=48, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb = _text_box(sl, Inches(1), Inches(4), Inches(11), Inches(1))
_set_text(tb.text_frame, "Questions?", size=28, colour=ACCENT, align=PP_ALIGN.CENTER)
tb = _text_box(sl, Inches(1), Inches(5.5), Inches(11), Inches(1))
_set_text(tb.text_frame,
    "github.com/partham20/S-Board-Firmware\n"
    "github.com/partham20/boot_manager\n"
    "github.com/partham20/fw-updater-gui",
    size=14, colour=GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
out = "docs/OTA_Firmware_Update_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
